#!/usr/bin/env python3
"""
Render PPTX slides to PNG images using python-pptx + Pillow.
Handles solid-fill backgrounds, text shapes with formatting, and embedded images.
"""
import os, re, sys, io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import zipfile

# Output resolution: 150 DPI equivalent at 13.33 x 7.5 inches
OUTPUT_W = 1600
OUTPUT_H = 900

EMU_W = 12192000
EMU_H = 6858000

def emu_to_px(emu, total_emu, total_px):
    return int(emu / total_emu * total_px)

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 6:
        return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))
    return (200, 200, 200)

def get_shape_fill(shape):
    """Extract fill color from a shape. Returns RGB tuple or None."""
    try:
        sp = shape.element
        # Check spPr solidFill
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            spPr = sp.find(qn('p:grpSpPr'))
        if spPr is not None:
            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    return hex_to_rgb(srgb.get('val', ''))
                schemeClr = solidFill.find(qn('a:schemeClr'))
                if schemeClr is not None:
                    return None  # theme color, skip
    except Exception:
        pass
    return None

def get_slide_bg(slide):
    """Get slide background color."""
    try:
        bg = slide.element.find(qn('p:cSld'))
        if bg is not None:
            bgPr = bg.find(qn('p:bg'))
            if bgPr is not None:
                bgPr2 = bgPr.find(qn('p:bgPr'))
                if bgPr2 is not None:
                    solidFill = bgPr2.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgb = solidFill.find(qn('a:srgbClr'))
                        if srgb is not None:
                            return hex_to_rgb(srgb.get('val', ''))
    except Exception:
        pass
    return (11, 33, 52)  # default dark blue

def get_text_color(run_elem):
    """Extract text color from run element."""
    try:
        rPr = run_elem.find(qn('a:rPr'))
        if rPr is not None:
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    return hex_to_rgb(srgb.get('val', ''))
    except Exception:
        pass
    return None

def get_font_size(run_elem):
    """Extract font size in points from run element."""
    try:
        rPr = run_elem.find(qn('a:rPr'))
        if rPr is not None:
            sz = rPr.get('sz')
            if sz:
                return int(sz) / 100  # hundredths of a point
    except Exception:
        pass
    return None

def is_bold(run_elem):
    try:
        rPr = run_elem.find(qn('a:rPr'))
        if rPr is not None:
            return rPr.get('b') in ('1', 'true')
    except Exception:
        pass
    return False

def load_font(size, bold=False):
    """Load a system font at given size."""
    font_paths = [
        '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf' if bold else
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf' if bold else
        '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
        '/usr/share/fonts/truetype/freefont/FreeSansBold.ttf' if bold else
        '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
    ]
    for path in font_paths:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, int(size))
            except Exception:
                pass
    return ImageFont.load_default()

def wrap_text(text, font, max_width, draw):
    """Wrap text to fit within max_width pixels."""
    words = text.split()
    if not words:
        return ['']
    lines = []
    current = []
    for word in words:
        test = ' '.join(current + [word])
        bbox = draw.textbbox((0, 0), test, font=font)
        w = bbox[2] - bbox[0]
        if w <= max_width or not current:
            current.append(word)
        else:
            lines.append(' '.join(current))
            current = [word]
    if current:
        lines.append(' '.join(current))
    return lines

def get_para_align(para_elem):
    """Get paragraph alignment."""
    pPr = para_elem.find(qn('a:pPr'))
    if pPr is not None:
        algn = pPr.get('algn')
        if algn == 'ctr':
            return 'center'
        elif algn == 'r':
            return 'right'
    return 'left'

def extract_embedded_images(pptx_path):
    """Extract all images from pptx zip."""
    images = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if re.match(r'ppt/media/.*\.(png|jpg|jpeg|gif)', name, re.I):
                images[name] = z.read(name)
    return images

def get_pic_media_ref(shape_elem, slide):
    """Get media reference for a picture shape."""
    try:
        blip = shape_elem.find('.//' + qn('a:blip'))
        if blip is not None:
            rEmbed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rEmbed and hasattr(slide, 'part'):
                rel = slide.part.rels.get(rEmbed)
                if rel:
                    return rel.target_ref
    except Exception:
        pass
    return None

def render_slide(slide, pptx_path, slide_idx):
    """Render a single slide to a PIL Image."""
    img = Image.new('RGB', (OUTPUT_W, OUTPUT_H), get_slide_bg(slide))
    draw = ImageDraw.Draw(img)

    # Load embedded images once
    embedded = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if re.match(r'ppt/media/', name) and re.search(r'\.(png|jpg|jpeg)$', name, re.I):
                try:
                    embedded[os.path.basename(name)] = Image.open(io.BytesIO(z.read(name))).convert('RGBA')
                except Exception:
                    pass

    for shape in slide.shapes:
        if not shape.has_text_frame and shape.shape_type == 13:  # Picture
            try:
                left = emu_to_px(shape.left or 0, EMU_W, OUTPUT_W)
                top = emu_to_px(shape.top or 0, EMU_H, OUTPUT_H)
                width = emu_to_px(shape.width or 0, EMU_W, OUTPUT_W)
                height = emu_to_px(shape.height or 0, EMU_H, OUTPUT_H)
                # Try to find embedded image
                rEmbed = None
                blip = shape.element.find('.//' + qn('a:blip'))
                if blip is not None:
                    rEmbed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if rEmbed and hasattr(slide, 'part'):
                    rel = slide.part.rels.get(rEmbed)
                    if rel:
                        target = rel.target_ref
                        fname = os.path.basename(target)
                        if fname in embedded and width > 0 and height > 0:
                            pic = embedded[fname].resize((width, height), Image.LANCZOS)
                            img.paste(pic, (left, top), pic.split()[3] if pic.mode == 'RGBA' else None)
            except Exception:
                pass

        # Draw filled shapes (rectangles, ellipses)
        fill = get_shape_fill(shape)
        if fill and not shape.has_text_frame:
            try:
                left = emu_to_px(shape.left or 0, EMU_W, OUTPUT_W)
                top = emu_to_px(shape.top or 0, EMU_H, OUTPUT_H)
                width = emu_to_px(shape.width or 0, EMU_W, OUTPUT_W)
                height = emu_to_px(shape.height or 0, EMU_H, OUTPUT_H)
                # Get alpha
                alpha_val = 255
                sp = shape.element
                spPr = sp.find(qn('p:spPr'))
                if spPr is not None:
                    solidFill = spPr.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgb = solidFill.find(qn('a:srgbClr'))
                        if srgb is not None:
                            alpha_elem = srgb.find(qn('a:alpha'))
                            if alpha_elem is not None:
                                alpha_pct = int(alpha_elem.get('val', '100000')) / 100000
                                alpha_val = int(alpha_pct * 255)
                color_rgba = fill + (alpha_val,)
                prstGeom = None
                spPr2 = shape.element.find(qn('p:spPr'))
                if spPr2 is not None:
                    prstGeom = spPr2.find(qn('a:prstGeom'))
                is_ellipse = prstGeom is not None and prstGeom.get('prst') == 'ellipse'
                overlay = Image.new('RGBA', (OUTPUT_W, OUTPUT_H), (0,0,0,0))
                od = ImageDraw.Draw(overlay)
                if is_ellipse:
                    od.ellipse([left, top, left+width, top+height], fill=color_rgba)
                else:
                    od.rectangle([left, top, left+width, top+height], fill=color_rgba)
                img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                draw = ImageDraw.Draw(img)
            except Exception:
                pass

        # Render text frames
        if shape.has_text_frame:
            try:
                left = emu_to_px(shape.left or 0, EMU_W, OUTPUT_W)
                top = emu_to_px(shape.top or 0, EMU_H, OUTPUT_H)
                width = emu_to_px(shape.width or 1, EMU_W, OUTPUT_W)
                height = emu_to_px(shape.height or 1, EMU_H, OUTPUT_H)

                # Background fill for text shape
                fill = get_shape_fill(shape)
                if fill:
                    draw.rectangle([left, top, left+width, top+height], fill=fill)

                y_cursor = top + 8
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        y_cursor += 4
                        continue

                    # Determine font size from first run
                    font_size = 14
                    run_bold = False
                    run_color = (231, 238, 243)  # default --ink

                    runs = para._p.findall(qn('a:r'))
                    if runs:
                        fsz = get_font_size(runs[0])
                        if fsz:
                            font_size = max(8, min(80, fsz))
                        run_bold = is_bold(runs[0])
                        rc = get_text_color(runs[0])
                        if rc:
                            run_color = rc

                    # Scale font size to output resolution
                    scaled_size = font_size * (OUTPUT_W / (13.33 * 96))
                    scaled_size = max(8, min(72, scaled_size))

                    font = load_font(scaled_size, run_bold)
                    align = get_para_align(para._p)

                    # Word wrap
                    margin = 8
                    available_w = max(50, width - margin * 2)
                    lines = wrap_text(para_text, font, available_w, draw)

                    for line in lines:
                        if y_cursor > top + height:
                            break
                        if align == 'center':
                            bbox = draw.textbbox((0, 0), line, font=font)
                            tw = bbox[2] - bbox[0]
                            x = left + (width - tw) // 2
                        elif align == 'right':
                            bbox = draw.textbbox((0, 0), line, font=font)
                            tw = bbox[2] - bbox[0]
                            x = left + width - tw - margin
                        else:
                            x = left + margin

                        # Shadow for readability
                        draw.text((x+1, y_cursor+1), line, font=font, fill=(0,0,0,128))
                        draw.text((x, y_cursor), line, font=font, fill=run_color)

                        bbox = draw.textbbox((x, y_cursor), line, font=font)
                        line_h = bbox[3] - bbox[1]
                        y_cursor += max(line_h + 3, int(scaled_size * 1.2))

                    # Para spacing
                    y_cursor += 4

            except Exception as e:
                pass

    return img

def render_deck(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        img = render_slide(slide, pptx_path, i)
        out_path = os.path.join(output_dir, f'slide-{i+1:02d}.png')
        img.save(out_path, 'PNG', optimize=True)
        print(f"  Rendered slide {i+1}/{total}: {out_path}")
    print(f"Done: {total} slides in {output_dir}")

if __name__ == '__main__':
    base = '/home/user/MRA-Files/kidney-care'
    print("Rendering Deck 1...")
    render_deck(
        f'{base}/assets/source/Deck1_Independent_JV.pptx',
        f'{base}/assets/deck1'
    )
    print("\nRendering Deck 2...")
    render_deck(
        f'{base}/assets/source/Deck2_Chain_JV.pptx',
        f'{base}/assets/deck2'
    )
